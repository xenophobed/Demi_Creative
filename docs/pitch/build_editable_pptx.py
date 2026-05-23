#!/usr/bin/env python3
"""
Build an EDITABLE PowerPoint of the Kids Creative Workshop pitch deck.

Unlike keynote-deck.pptx (compiled by Marp, where every slide is a flat
rasterised image), this output has real editable text boxes, tables, and
movable images — open it in PowerPoint or Keynote and edit any word.

SVG diagrams are embedded as PNG pictures (converted via macOS qlmanage):
they stay movable/resizable but are not text-editable. Everything else is.

Run:  python docs/pitch/build_editable_pptx.py
Out:  docs/pitch/keynote-deck-editable.pptx
"""

import re
import subprocess
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

PITCH = Path(__file__).parent
ASSETS = PITCH / "assets"
TMP = PITCH / ".pptx-build"
TMP.mkdir(exist_ok=True)

# ── Brand palette ───────────────────────────────────────────────────────
PINK = RGBColor(0xDB, 0x27, 0x77)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
LIGHT_TINT = RGBColor(0xFD, 0xF2, 0xF8)   # title-card background
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
PINK_DARK = RGBColor(0xF4, 0x72, 0xB6)    # pink on dark slides
LIGHT_TEXT = RGBColor(0xF8, 0xFA, 0xFC)

FONT = "Helvetica Neue"
MONO = "Menlo"

# ── SVG → PNG (qlmanage + transparent-margin crop) ──────────────────────
def svg_to_png(svg_name: str) -> Path:
    svg = ASSETS / svg_name
    out = TMP / (svg.stem + ".png")
    subprocess.run(
        ["qlmanage", "-t", "-s", "2400", "-o", str(TMP), str(svg)],
        capture_output=True, check=True,
    )
    raw = TMP / (svg.name + ".png")          # qlmanage names it <file>.svg.png
    im = Image.open(raw).convert("RGBA")
    bbox = im.getbbox()                       # crop transparent padding
    if bbox:
        im = im.crop(bbox)
    im.save(out)
    raw.unlink(missing_ok=True)
    return out

# ── Speaker notes: parse them out of keynote-deck.md ────────────────────
def load_notes() -> dict[int, str]:
    text = (PITCH / "keynote-deck.md").read_text()
    notes: dict[int, str] = {}
    for m in re.finditer(
        r"🎤 SCRIPT · Slide (\d+).*?\n(.*?)(?=\n-->)", text, re.DOTALL
    ):
        notes[int(m.group(1))] = m.group(2).strip()
    return notes

NOTES = load_notes()

# ── Presentation scaffold (16:9 widescreen) ─────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def new_slide(bg: RGBColor = WHITE):
    slide = prs.slides.add_slide(BLANK)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def set_notes(slide, n: int):
    if n in NOTES:
        slide.notes_slide.notes_text_frame.text = NOTES[n]


_MD = re.compile(r"(\*\*.+?\*\*|`.+?`|\*.+?\*)")


def _runs(paragraph, text, size, color, base_bold=False):
    """Render markdown-ish text (**bold**, `code`, *italic*) as pptx runs."""
    for chunk in _MD.split(text):
        if not chunk:
            continue
        run = paragraph.add_run()
        bold, italic, mono = base_bold, False, False
        if chunk.startswith("**") and chunk.endswith("**"):
            chunk, bold = chunk[2:-2], True
        elif chunk.startswith("`") and chunk.endswith("`"):
            chunk, mono = chunk[1:-1], True
        elif chunk.startswith("*") and chunk.endswith("*"):
            chunk, italic = chunk[1:-1], True
        run.text = chunk
        f = run.font
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.italic = italic
        f.name = MONO if mono else FONT


def textbox(slide, left, top, width, height, lines, *, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, bold) tuples."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        _runs(p, text, size, color, base_bold=bold)
    return box


def heading(slide, text, *, color=SLATE):
    return textbox(
        slide, Inches(0.7), Inches(0.45), Inches(12.0), Inches(1.0),
        [(text, 32, color, True)],
    )


def caption(slide, text, *, color=MUTED, top=6.7):
    return textbox(
        slide, Inches(0.7), Inches(top), Inches(12.0), Inches(0.6),
        [(text, 15, color, False)],
    )


def add_table(slide, rows, *, left=0.7, top=1.6, width=11.93, height=4.6,
              header=True, col_widths=None):
    n_rows, n_cols = len(rows), len(rows[0])
    gt = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    table = gt.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            is_header = header and r == 0
            cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF3, 0xFF) if is_header else WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            size = 16 if is_header else 14
            color = VIOLET if is_header else SLATE
            _runs(p, cell_text, size, color, base_bold=is_header)
    return table


def add_image(slide, png: Path, *, max_w=11.5, max_h=4.7, top=1.7):
    im = Image.open(png)
    iw, ih = im.size
    ratio = iw / ih
    w, h = max_w, max_w / ratio
    if h > max_h:
        h, w = max_h, max_h * ratio
    left = (13.333 - w) / 2
    slide.shapes.add_picture(
        str(png), Inches(left), Inches(top), Inches(w), Inches(h)
    )


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
s = new_slide(LIGHT_TINT)
textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.4),
        [("Kids Creative Workshop", 54, PINK, True)])
textbox(s, Inches(1.0), Inches(3.6), Inches(11.3), Inches(1.0),
        [("An *agentic* app for kids — built on Claude Agent SDK.", 30, VIOLET, True)])
textbox(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.6),
        [("6-minute pitch · Single agent → agent team · 2026", 17, MUTED, False)])
set_notes(s, 1)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Three moments
# ════════════════════════════════════════════════════════════════════════
s = new_slide(LIGHT_TINT)
textbox(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(1.1),
        [("Three moments. One buddy.", 44, PINK, True)])
textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(3.6), [
    ("Her drawing — becomes a story.", 34, VIOLET, True),
    ("Her character — stars in the next adventure.", 34, VIOLET, True),
    ("Tomorrow's news — arrives as her podcast.", 34, VIOLET, True),
])
set_notes(s, 2)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Four problems
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Today's AI fails kids on four fronts")
add_table(s, [
    ["🚫 Not personalized enough", "🚫 Not highly customized"],
    ["Generic output. No memory of her characters, her voice, her age.",
     "Same prompt, same answer for every child. No buddy identity. No per-child capability set."],
    ["🚫 No suited news for kids", "🚫 No long-term persistence"],
    ["Today's news products are adult-first. No age-aware filter. No narrative voice for kids.",
     "Each session resets. Lightning the puppy is forgotten by Monday."],
], header=False, top=1.7, height=4.2)
caption(s, "Existing AI extracts. We collaborate.")
set_notes(s, 3)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Our agent's abilities
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Meet My Agent — one ability per problem")
add_table(s, [
    ["Problem", "Our agent's ability"],
    ["Not personalized enough",
     "**Persona + character memory** — buddy is named & customized; recurring characters recalled across sessions (`character_repo`)"],
    ["Not highly customized",
     "**Per-child `AgentDefinition` + skills gating** — age-aware capabilities (3–5 / 6–8 / 9–12); persona shared as system context to every specialist"],
    ["No suited news for kids",
     "**`kids_daily` specialist** — age-stratified prompts + per-reply safety review; news arrives as a kid-safe podcast in the buddy's voice"],
    ["No long-term persistence",
     "**`agent_repo` + `character_repo` + vector search** — buddy persona and recurring characters survive every session; one buddy, for life"],
], top=1.6, height=4.6, col_widths=[3.2, 8.73])
caption(s, "What kids feel: interactive (streaming) · proactive (recommendations + recall) · persistent (memory).")
set_notes(s, 4)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Six agentic features
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Foundation — six agentic features, one stack")
add_table(s, [
    ["Agentic feature", "How we build it"],
    ["🌊 Interactive", "Streaming **SSE** · async-generator agents · live tool-use events"],
    ["🎯 Responsive", "**LLM model** (Claude Haiku for speed) · deterministic + LLM intent routing · per-agent skill curation"],
    ["💡 Proactive", "**Prompt engineering** (system-prompt scaffolding) · **vector DB** recall · `character_repo` lookups"],
    ["🧠 Persistent", "**Vector DB** (ChromaDB / pgvector) · `agent_repo` + `character_repo` (SQL) · cross-session memory"],
    ["🛡️ Reactive", "**MCP** safety tool · `@tool` decorator · `enforce_chat_safety` + suggest-and-retry · age-aware thresholds (0.90 / 0.85)"],
    ["🚀 Autonomous (future)", "Multi-step planning · **skills** composition · self-prompted explore loops · scheduled buddy initiatives"],
], top=1.5, height=4.9, col_widths=[3.0, 8.93])
caption(s, "Built from: prompt engineering · MCP · tools · skills · LLM model · vector DB.")
set_notes(s, 5)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Four architecture patterns
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Four agent architecture patterns — we use all four")
add_table(s, [
    ["#", "Pattern", "What it does", "Where we use it"],
    ["1", "🤖 Single agent", "One agent, one job · linear inference", "Straight TTS via `audio_narration`"],
    ["2", "🔀 Sub-agent fan-out", "Same task spawned in parallel for speed", "Concurrent vision crops · parallel `character_repo` lookups"],
    ["3", "👥 Agent team", "Multiple agents collaborate by **role** · via `AgentDefinition`", "**My Agent**: proxy + 4 role specialists + `safety_review`"],
    ["4", "🎼 Multi-agent orchestrator", "Agents created **dynamically** · A2A extensible", "Proxy registers new `AgentDefinition`s at runtime"],
], top=1.6, height=4.4, col_widths=[0.5, 2.7, 4.3, 4.43])
caption(s, "Shared state flows through build_my_agent_context(). A2A extends to external teams (future).")
set_notes(s, 6)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Six memory types
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Six memory types — one buddy, layered recall")
add_table(s, [
    ["#", "Memory type", "What it stores", "Backed by"],
    ["1", "🗨️ Session", "This-chat conversation history", "`agent_chat_repository`"],
    ["2", "⚡ Working", "Per-turn execution context (tool results, persona)", "`build_my_agent_context()`"],
    ["3", "📅 Episodic", "Past creations — stories, podcasts, choices", "`stories` · `interactive_sessions` · `kids_daily_episodes`"],
    ["4", "📋 Factual", "Buddy persona, child profile, preferences", "`agent_repo` · `preference_repository` · `users`"],
    ["5", "🧠 Semantic", "Embeddings of characters, themes, style", "ChromaDB / pgvector via `vector_search_server`"],
    ["6", "🛠️ Procedural", "How to generate each content type", "`prompts/*.md` · `@tool` skills · `enabled_skills`"],
], top=1.5, height=4.9, col_widths=[0.5, 2.4, 4.4, 4.63])
caption(s, "The buddy remembers, understands, acts, talks, and reasons in flight.")
set_notes(s, 7)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — The team (dark, architecture diagram)
# ════════════════════════════════════════════════════════════════════════
s = new_slide(DARK_BG)
textbox(s, Inches(0.7), Inches(0.4), Inches(12.0), Inches(0.9),
        [("The team — one proxy, four specialists, one safety gate", 30, PINK_DARK, True)])
add_image(s, svg_to_png("architecture.svg"), max_w=11.0, max_h=5.0, top=1.5)
textbox(s, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.5),
        [("Built on Claude Agent SDK + custom MCP tool servers.", 14, RGBColor(0x94, 0xA3, 0xB8), False)])
set_notes(s, 8)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Three-layer infrastructure
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Three-layer infrastructure — each service does one thing")
add_image(s, svg_to_png("infrastructure.svg"), max_w=11.5, max_h=5.1, top=1.5)
set_notes(s, 9)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Backend layered architecture
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Backend — seven layers, one direction")
add_image(s, svg_to_png("backend-layers.svg"), max_w=10.5, max_h=5.3, top=1.4)
set_notes(s, 10)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 11 — The buddy (3-state strip)
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "The buddy — three states, one identity")
add_image(s, ASSETS / "buddy-states.png", max_w=12.0, max_h=4.6, top=1.6)
caption(s, "Empty state → Customize (name, avatar, theme) → Chat. Three React states, one persona, "
           "persisted across every session.")
set_notes(s, 11)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 12 — What the buddy creates
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "What the buddy creates — three surfaces, one character")
add_image(s, ASSETS / "system-creates.png", max_w=12.2, max_h=3.6, top=1.45)
add_table(s, [
    ["📖 Image-to-Story", "🌟 Interactive Story", "🎙️ Kids Daily"],
    ["The Singing Shells of Coral Bay — art story with illustrated cover",
     "Ember and the Golden Dragon — branching scene with 3 choice cards",
     "Amazing Animals and How We Keep Them Safe — episode + transcript + audio"],
], header=True, top=5.25, height=1.1)
caption(s, "Same character — Ember + the recurring crew — across all three. Real outputs from the live app.",
        top=6.55)
set_notes(s, 12)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Community & sharing
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Community & sharing — COPPA by schema, not by policy")
add_table(s, [
    ["Where most products fail", "What we do"],
    ["Posts JOIN to `users.name` for byline",
     "`hub_posts.agent_name` is a **snapshot column** — written at post time, never JOINed"],
    ["`users.email` accidentally leaks via API",
     "Read paths can't reach `users` at all — schema doesn't allow it"],
    ["Safety is a code-review checklist",
     "Safety is a CHECK constraint + invariant test"],
], top=1.7, height=3.6, col_widths=[5.0, 6.93])
caption(s, "Every Hub post is bylined by the buddy persona, never the child. "
           "Verified by test_hub_coppa_invariant.py.")
set_notes(s, 13)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Open by design (code block)
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Open by design — one AgentDefinition adds a specialist")
code = (
    '# Adding a "music_story" specialist to the agent team:\n'
    'proxy.register(AgentDefinition(\n'
    '    name="music_story",\n'
    '    model="haiku",\n'
    '    system_prompt=Path("prompts/music-story.md").read_text(),\n'
    '    tools=["music_generator", "vector_search"],\n'
    '    enabled_skills=["compose"],\n'
    '))\n\n'
    '# Routing picks it up automatically. Safety gate runs on every reply.\n'
    '# Shared context (persona, child_id, recurring chars) flows in.'
)
box = s.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.3), Inches(3.6))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
box.line.color.rgb = CARD_BORDER
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = Inches(0.3)
tf.margin_top = tf.margin_bottom = Inches(0.2)
for i, line in enumerate(code.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = line or " "
    r.font.name = MONO
    r.font.size = Pt(15)
    r.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A) if not line.startswith("#") else MUTED
caption(s, "A2A bridge (future) extends to external agent teams via the same registration contract.",
        top=5.8)
set_notes(s, 14)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Roadmap
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Roadmap — two phases shipped, two ahead")
add_image(s, svg_to_png("roadmap.svg"), max_w=12.0, max_h=4.8, top=1.7)
set_notes(s, 15)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Where we are
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Where we are")
add_table(s, [
    ["Milestone", "Status"],
    ["Phase 1 — MVP: single agent + image-to-story + safety + TTS", "✅ 92/92 shipped"],
    ["Phase 2 — multi-agent team + memory + news + community", "✅ 180/180 shipped"],
    ["Phase 3 — video · parent dashboard · gamification", "🔜 In design"],
], top=1.7, height=2.6, col_widths=[9.0, 2.93])
textbox(s, Inches(0.7), Inches(4.7), Inches(12.0), Inches(1.4), [
    ("Engineering rigor: 700+ contract tests · per-reply programmatic safety (age-aware) · "
     "silent safety-bypass caught + fixed in 24h · merge-train of 7 PRs landed last week.",
     18, SLATE, False),
])
caption(s, "Add real numbers before presenting: pilot users · sessions/week · feedback quotes.",
        top=6.2)
set_notes(s, 16)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Failures we owned
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Failures we owned — receipts, not theater")
add_table(s, [
    ["What we tried", "How it broke", "What we did about it"],
    ["**SDK subprocess** for image-to-story",
     "Railway exit -9 · OOM kills under load",
     "Ported all 3 generation agents to direct API · ~50% memory drop"],
    ["`await check_content_safety({...})`",
     "`SdkMcpTool` wrapper not callable · `TypeError` swallowed → default 0.9 score",
     "Caught + fixed in 24h · `.handler` calling convention · 3 agents, 1 PR"],
    ["**Single agent + safety prompt**",
     "Model occasionally produced unsafe replies",
     "Per-reply programmatic safety subagent · age-aware · fail-closed retry"],
], top=1.7, height=4.0, col_widths=[3.3, 4.3, 4.33])
caption(s, "Most pitches hide bugs. We name ours — that's how you know we run safety like infrastructure.")
set_notes(s, 17)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Why this matters (closing)
# ════════════════════════════════════════════════════════════════════════
s = new_slide(LIGHT_TINT)
textbox(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0),
        [("Why this matters", 44, PINK, True)])
textbox(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(3.2), [
    ("Agentic from day one — not a wrapper, not a prompt. Real SDK, real tools, real orchestration.", 21, SLATE, False),
    ("272 stories shipped across 3 milestones — execution proof.", 21, SLATE, False),
    ("Programmatic safety on every reply — non-negotiable, code-enforced, not vibes.", 21, SLATE, False),
    ("Community that protects child PII at the schema level — COPPA by construction.", 21, SLATE, False),
    ("A buddy that grows with the child — character continuity across image, story, podcast, share.", 21, SLATE, False),
])
textbox(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.0),
        [("AI that grows up *with* kids — safely.", 36, PINK, True)])
set_notes(s, 18)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Appendix
# ════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Appendix — technical deep-dive (backup for Q&A)")
add_table(s, [
    ["Topic", "One-line answer"],
    ["Agent", "`AgentDefinition(model, system_prompt, tools, enabled_skills)` — one specialist w/ a curated capability set"],
    ["Subagent", "An agent registered under the proxy's `agents=` dict · invoked via the SDK's Agent tool delegation"],
    ["Agent team", "Proxy + 3 product subagents + audio_narration tool + safety_review · all share the context bus"],
    ["Orchestrator", "The proxy — routes intent · composes specialist outputs · runs safety_review on every reply"],
    ["Why this shape", "Bigger prompt degrades w/ specialty count · prompt chaining = no shared state · agent team = shared context + parallel specialty"],
    ["Per-reply safety", "`enforce_chat_safety()` after every reply · age-aware threshold · suggest-improvements retry · `safety_blocked` telemetry"],
    ["COPPA pattern", "`hub_posts.agent_name/avatar/title` — immutable snapshot columns; no read path JOINs `users`"],
    ["Tech stack", "FastAPI + Pydantic v2 · SQLite (dev) / Postgres + pgvector (prod) · React 18 + TS + Tailwind"],
], top=1.5, height=5.0, col_widths=[2.4, 9.53])
set_notes(s, 19)

# ── Save ────────────────────────────────────────────────────────────────
out = PITCH / "keynote-deck-editable.pptx"
prs.save(str(out))
print(f"Wrote {out} — {len(prs.slides)} slides, fully editable.")
